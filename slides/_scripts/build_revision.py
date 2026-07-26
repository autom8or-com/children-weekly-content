"""Standalone builder for the once-a-term REVISION deck + flip-card quiz decks.

The revision format (Pre-K "Who is this? / What plague is this?" flip-card quiz)
is intentionally NOT wired into the core lesson-deck engine (deck_render.py /
deck_html.py / build_pptx.py). It is a rare deliverable, so instead of leaving
flip-card slide types in the shared renderers forever (and keeping the native +
screenshot renderers in lockstep for them), this script *imports* the visual
primitives from deck_render and defines only the new things it needs.
Native / editable output only.

Schema (slides/projects/<project>/slides.json is a list of typed dicts):
  title         {eyebrow, title, subtitle?, footer?}          -> render_title_goodbye
  goodbye       {eyebrow, title, subtitle?, footer?}          -> render_title_goodbye
  section       {palette, eyebrow, title, footer}             -> render_section_header
  qa-question   {palette, question, hint?, prompt?, image}    -> render_card (card + image)
  qa-answer     {palette, chip?, answer, lesson?}             -> render_qa_answer (big box)
  num-question  {palette, number, text, footnote?}            -> render_num_card (2 cards)
  num-answer    {palette, number, text, footnote?, chip?}     -> render_num_card (2 cards)
  diary-list    {title, lead?, items:[{emoji, story, lesson}]}-> render_diary_list

`image` (qa-question only) is a path relative to slides/.

The num-question / num-answer types are used by flip-card quiz decks (e.g. the
"Graduation Games" Pre-Teens John 7 deck) where each content slide has TWO cards
side by side: a number card on the left and a text card on the right. The Q-A
pair shares the same number — so a kid can flip from the question to the answer
and visually pair them by the number.

Run:  python3 slides/_scripts/build_revision.py [--project NAME] [--out FILE.pptx]
       (default --project=revision -> projects/revision/, default --out=Revision.pptx)
"""
import argparse
import json
import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

import deck_render as dr
from deck_render import (
    render_title_goodbye, render_section_header, render_card,
    render_prayer, render_house_rules,
    _rect, _bg, _textbox, PALETTES, TEXT_DARK, TEXT_BROWN, BG_CREAM,
    FONT_LABEL, FONT_HEAD, IN_PER_PX,
)

SLIDES_DIR = Path(__file__).resolve().parent.parent          # .../slides
PARTIALS = SLIDES_DIR / "_templates" / "partials"
SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)


def _project_dir(name: str) -> Path:
    """Resolve the project directory. Defaults to 'revision' for back-compat."""
    p = SLIDES_DIR / "projects" / name
    if not p.exists():
        raise SystemExit(
            f"[ERROR] project directory not found: {p}\n"
            f"        Pass --project <name> for a project under slides/projects/."
        )
    return p


def resolve_image(ref):
    """slides.json image ref (path relative to slides/) -> absolute path str or None."""
    if not ref:
        return None
    p = (SLIDES_DIR / ref).resolve()
    return str(p) if p.exists() else None


# ─────────────────────────── num-card layout (Pre-Teens flip quiz) ───────────────────────────

# Stage is 1920x1080 px (scaled via IN_PER_PX). The two cards sit side by side.
# Number card on the left, content card on the right; identical Y/H so they line up.
_NUM_CARD_TOP = 140
_NUM_CARD_H = 800
_NUM_CARD_GAP = 40
_NUM_CARD_LEFT_X = 80
_NUM_CARD_LEFT_W = 460
_NUM_CARD_RIGHT_X = _NUM_CARD_LEFT_X + _NUM_CARD_LEFT_W + _NUM_CARD_GAP  # 580
_NUM_CARD_RIGHT_W = 1920 - _NUM_CARD_RIGHT_X - 80                       # 1260


def render_num_card(prs, s, *, mode):
    """Two-card flip-quiz layout (number card + content card).

    mode: "question" or "answer" — controls chip text and font size tuning.
           "question" makes the content card the bigger headline (the prompt).
           "answer" makes the chip ("✅ ANSWER") show above the answer.
    """
    pal = PALETTES.get(s.get("palette", "preamble"), PALETTES["preamble"])
    number = str(s["number"])
    text = s["text"]
    footnote = s.get("footnote")
    chip_text = s.get("chip") or ("❓ QUESTION" if mode == "question" else "✅ ANSWER")

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide, pal["sec_bg"])

    # ── LEFT: number card (accent bg, dark text — the constant visual anchor) ──
    _rect(slide, _NUM_CARD_LEFT_X, _NUM_CARD_TOP, _NUM_CARD_LEFT_W, _NUM_CARD_H,
          pal["accent"], rounded=True)
    # Big number, vertically centered. Caveat has a tall ascent+descent, so we
    # cap by character count: 1 digit gets 280pt, 2 digits get 200pt, 3+ → 150pt.
    n_chars = len(number)
    if n_chars <= 1:
        num_size = 280
    elif n_chars == 2:
        num_size = 200
    else:
        num_size = 150
    _textbox(
        slide,
        _NUM_CARD_LEFT_X * IN_PER_PX,
        _NUM_CARD_TOP * IN_PER_PX,
        _NUM_CARD_LEFT_W * IN_PER_PX,
        _NUM_CARD_H * IN_PER_PX,
        [{"text": number, "font": FONT_HEAD, "size": num_size,
          "color": BG_CREAM, "bold": True, "space_after": 0, "line_spacing": 1.0}],
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, auto_fit_shrink=False,
    )

    # ── RIGHT: content card (card bg, dark text — the variable part) ──
    _rect(slide, _NUM_CARD_RIGHT_X, _NUM_CARD_TOP, _NUM_CARD_RIGHT_W, _NUM_CARD_H,
          pal["card"], rounded=True)
    # accent cap stripe across the top (matches the qa-answer style)
    _rect(slide, _NUM_CARD_RIGHT_X, _NUM_CARD_TOP, _NUM_CARD_RIGHT_W, 14, pal["accent"],
          rounded=True)

    # Compose content: chip → text → optional footnote. Bigger font for the
    # answer (the reveal) than the question (the prompt).
    title_size = 44 if mode == "question" else 54
    lines = [
        {"text": chip_text, "font": FONT_LABEL, "size": 28, "color": pal["accent"],
         "space_after": 16, "line_spacing": 1.0},
        {"text": text, "size": title_size, "color": TEXT_DARK, "bold": True,
         "space_after": 0, "line_spacing": 1.2},
    ]
    if footnote:
        lines.append({"text": "▬▬▬", "font": FONT_LABEL, "size": 16,
                      "color": pal["accent"], "space_after": 10, "line_spacing": 1.0})
        lines.append({"text": footnote, "size": 28, "color": TEXT_BROWN,
                      "space_after": 0, "line_spacing": 1.2})
    _textbox(
        slide,
        (_NUM_CARD_RIGHT_X + 50) * IN_PER_PX,
        (_NUM_CARD_TOP + 50) * IN_PER_PX,
        (_NUM_CARD_RIGHT_W - 100) * IN_PER_PX,
        (_NUM_CARD_H - 100) * IN_PER_PX,
        lines,
        align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE, auto_fit_shrink=True,
    )
    return slide


# ─────────────────────────── original revision: qa-answer reveal ───────────────────────────

def render_qa_answer(prs, s):
    """Big-box answer reveal — one large card, '✅ ANSWER' chip, huge answer word,
    optional one-line Pre-K lesson. Distinct from the two-box question slide."""
    pal = PALETTES.get(s.get("palette", "preamble"), PALETTES["preamble"])
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide, pal["sec_bg"])

    cx, cy, cw, ch = 170, 150, 1580, 780
    _rect(slide, cx, cy, cw, ch, pal["card"], rounded=True)
    _rect(slide, cx, cy, cw, 14, pal["accent"], rounded=True)   # accent cap stripe

    lines = [
        {"text": s.get("chip", "✅ ANSWER"), "font": FONT_LABEL, "size": 40,
         "color": pal["accent"], "space_after": 22, "line_spacing": 1.0},
        {"text": s["answer"], "size": 120, "color": TEXT_DARK, "bold": True,
         "space_after": 18, "line_spacing": 0.95},
    ]
    if s.get("lesson"):
        lines.append({"text": "▬▬▬", "font": FONT_LABEL, "size": 20,
                      "color": pal["accent"], "space_after": 18, "line_spacing": 1.0})
        lines.append({"text": s["lesson"], "size": 48, "color": TEXT_BROWN,
                      "space_after": 0, "line_spacing": 1.05})
    _textbox(slide, (cx + 70) * IN_PER_PX, cy * IN_PER_PX, (cw - 140) * IN_PER_PX,
             ch * IN_PER_PX, lines, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
             auto_fit_shrink=True)
    return slide


def render_diary_list(prs, s):
    """Angela's-Diary roundup — one lesson per story, stacked rows (diary palette)."""
    pal = PALETTES["diary"]
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide, pal["sec_bg"])

    head = [
        {"text": s.get("chip", "📔 Angela's Diary"), "font": FONT_LABEL, "size": 32,
         "color": pal["accent"], "space_after": 14, "line_spacing": 1.0},
        {"text": s.get("title", "What We Learned"), "size": 60, "color": pal["hdr_title"],
         "bold": True, "space_after": 0, "line_spacing": 1.0},
    ]
    _textbox(slide, 200 * IN_PER_PX, 60 * IN_PER_PX, 1520 * IN_PER_PX, 210 * IN_PER_PX,
             head, align=PP_ALIGN.CENTER, auto_fit_shrink=False)

    items = s.get("items", [])
    n = len(items)
    if n == 0:
        return slide
    row_h, gap, y0 = 116, 18, 360
    for i, it in enumerate(items):
        ry = y0 + i * (row_h + gap)
        _rect(slide, 260, ry, 1400, row_h, pal["card"], rounded=True)
        _rect(slide, 260, ry, 10, row_h, pal["accent"])
        _textbox(slide, 300 * IN_PER_PX, (ry + 20) * IN_PER_PX, 110 * IN_PER_PX,
                 (row_h - 40) * IN_PER_PX,
                 [{"text": it.get("emoji", "•"), "size": 48, "color": TEXT_DARK,
                   "space_after": 0}], align=PP_ALIGN.CENTER, auto_fit_shrink=False)
        _textbox(slide, 430 * IN_PER_PX, (ry + 18) * IN_PER_PX, 1200 * IN_PER_PX,
                 (row_h - 34) * IN_PER_PX,
                 [{"text": f"**{it['story']}** — {it['lesson']}", "size": 30,
                   "color": TEXT_DARK, "space_after": 0, "line_spacing": 1.15}],
                 anchor=MSO_ANCHOR.MIDDLE, auto_fit_shrink=True)
    return slide


# ─────────────────────────── dispatch + build ───────────────────────────

def expand_includes(slides, partials_dir):
    """Replace {"include": "welcome"} with slides/_templates/partials/welcome.json.
    Extra keys on the include dict override the partial's fields."""
    out = []
    for s in slides:
        if isinstance(s, dict) and "include" in s:
            partial = json.loads((partials_dir / f"{s['include']}.json").read_text())
            overrides = {k: v for k, v in s.items() if k != "include"}
            partial.update(overrides)
            out.append(partial)
        else:
            out.append(s)
    return out


def render(prs, s, project_partials):
    t = s.get("type")
    if t == "title":
        return render_title_goodbye(prs, s, is_goodbye=False)
    if t == "goodbye":
        return render_title_goodbye(prs, s, is_goodbye=True)
    if t == "section":
        return render_section_header(prs, s)
    if t == "prayer":
        return render_prayer(prs, s, resolve_image(s.get("image")))
    if t == "house-rules":
        return render_house_rules(prs, s)
    if t == "qa-question":
        pal = PALETTES.get(s.get("palette", "preamble"), PALETTES["preamble"])
        body = []
        if s.get("hint"):
            body.append({"text": s["hint"], "size": 30, "color": TEXT_BROWN,
                         "space_after": 4, "line_spacing": 1.25})
        return render_card(
            prs, sec_bg=pal["sec_bg"], card_bg=pal["card"], accent=pal["accent"],
            chip=s.get("chip", "❓ QUESTION"), chip_font=FONT_LABEL, chip_col=pal["accent"],
            title=s["question"], body_lines=body, footer=s.get("prompt", "🤔 Can you guess?"),
            footer_col=TEXT_BROWN, image_path=resolve_image(s.get("image")),
            body_size=30, chip_size=22, title_size=56)
    if t == "qa-answer":
        return render_qa_answer(prs, s)
    if t == "num-question":
        return render_num_card(prs, s, mode="question")
    if t == "num-answer":
        return render_num_card(prs, s, mode="answer")
    if t == "diary-list":
        return render_diary_list(prs, s)
    raise ValueError(f"unknown slide type: {t!r}")


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("--project", default="revision",
                    help="Project name under slides/projects/ (default: revision)")
    ap.add_argument("--out", default=None,
                    help="Output PPTX filename (default: <Project>.pptx)")
    args = ap.parse_args()

    project = _project_dir(args.project)
    out_name = args.out or f"{args.project.capitalize()}.pptx"
    # Per-project partials override the default (e.g. project-specific welcome).
    project_partials = project / "partials"
    partials_dir = project_partials if project_partials.is_dir() else PARTIALS

    slides = expand_includes(json.loads((project / "slides.json").read_text()),
                             partials_dir)
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    missing = []
    for s in slides:
        if s.get("type") == "qa-question" and not resolve_image(s.get("image")):
            missing.append(s.get("image"))
        render(prs, s, project_partials)
    if missing:
        print("[WARN] missing images (rendered as placeholder blocks):")
        for m in missing:
            print("   ", m)

    out = project / out_name
    prs.save(str(out))
    print(f"✅ {len(slides)} slides -> {out}")


if __name__ == "__main__":
    main()
