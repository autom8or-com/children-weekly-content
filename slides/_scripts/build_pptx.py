"""Assemble a project's slides into a PowerPoint deck.

Two deck flavours, auto-detected from slides.json:

1. TYPED lesson deck (entries have a "type" key) — the storybook lesson format.
   Each slide is rendered as standalone HTML (deck_html.py), screenshotted at
   1920x1080 via playwright-cli, and placed full-bleed into the .pptx. Pixel-
   perfect to the Claude Design mock; slides are flat images.

2. LEGACY image deck (no "type") — full-bleed scene image + optional CTA text
   slide. Preserves older projects (e.g. frogs-and-gnats).

Usage:
  python build_pptx.py <project>
  python build_pptx.py <project> --out my_name.pptx
"""
import sys
import json
import shutil
import socket
import zipfile
import tempfile
import functools
import threading
import subprocess
from pathlib import Path
from http.server import ThreadingHTTPServer, SimpleHTTPRequestHandler

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

import deck_html
import deck_render

SLIDES_DIR = Path(__file__).parent.parent
SLIDE_W = Inches(13.33)   # 16:9 widescreen
SLIDE_H = Inches(7.5)
PW_SESSION = "deckbuild"
RENDER_WAIT = 1.5         # seconds to let fonts/layout settle before screenshot


# ─────────────────────────── legacy image deck ───────────────────────────

def add_image_slide(prs, image_path: Path, theme: str):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.shapes.add_picture(str(image_path), Emu(0), Emu(0), SLIDE_W, SLIDE_H)
    tb = slide.shapes.add_textbox(Inches(0.3), Inches(6.8), Inches(8), Inches(0.5))
    p = tb.text_frame.paragraphs[0]
    p.text = theme
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)


def add_cta_slide(prs, cta_text: str, theme: str):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 255, 255)

    tb_top = slide.shapes.add_textbox(Inches(0), Inches(0.4), SLIDE_W, Inches(0.6))
    p_top = tb_top.text_frame.paragraphs[0]
    p_top.text = theme
    p_top.alignment = PP_ALIGN.CENTER
    p_top.font.size = Pt(18)
    p_top.font.bold = True
    p_top.font.color.rgb = RGBColor(100, 100, 100)

    tb = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(11.33), Inches(5))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, line in enumerate(cta_text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = RGBColor(30, 30, 30)


# ─────────────────────────── typed lesson deck ───────────────────────────

def _free_port() -> int:
    with socket.socket(socket.AF_INET, socket.SOCK_STREAM) as s:
        s.bind(("127.0.0.1", 0))
        return s.getsockname()[1]


def _pw(*args) -> subprocess.CompletedProcess:
    return subprocess.run(["playwright-cli", f"-s={PW_SESSION}", *args],
                          capture_output=True, text=True)


def render_typed_deck(slides_meta, resolve_image) -> list[Path]:
    """Render each typed slide to a 1920x1080 PNG. Returns ordered PNG paths."""
    work = Path(tempfile.mkdtemp(prefix="deckbuild-"))
    html_dir = work / "html"
    png_dir = work / "png"
    html_dir.mkdir()
    png_dir.mkdir()

    # 1. write standalone HTML per slide (images inlined as data URIs)
    for i, slide in enumerate(slides_meta, 1):
        s = dict(slide)
        if s.get("image"):
            resolved = resolve_image(s["image"])
            s["image"] = str(resolved) if resolved else None
            if resolved is None:
                print(f"[WARN]  slide {i} ({slide.get('image')}) — image not found, placeholder used")
        (html_dir / f"slide-{i:02d}.html").write_text(deck_html.render_slide_html(s))

    # 2. serve the html dir (playwright blocks file://)
    class _QuietHandler(SimpleHTTPRequestHandler):
        def log_message(self, *a, **k):
            pass
    port = _free_port()
    handler = functools.partial(_QuietHandler, directory=str(html_dir))
    httpd = ThreadingHTTPServer(("127.0.0.1", port), handler)
    threading.Thread(target=httpd.serve_forever, daemon=True).start()

    pngs: list[Path] = []
    try:
        # 3. drive a headless browser: one screenshot per slide
        _pw("close")  # clear any stale session
        _pw("open")
        _pw("resize", "1920", "1080")
        for i in range(1, len(slides_meta) + 1):
            png = png_dir / f"slide-{i:02d}.png"
            _pw("goto", f"http://127.0.0.1:{port}/slide-{i:02d}.html")
            subprocess.run(["sleep", str(RENDER_WAIT)])
            r = _pw("screenshot", f"--filename={png}")
            if not png.exists():
                raise RuntimeError(f"screenshot failed for slide {i}: {r.stderr or r.stdout}")
            pngs.append(png)
            print(f"[RENDER] slide-{i:02d} ({slides_meta[i-1].get('type')})")
    finally:
        _pw("close")
        httpd.shutdown()
    return pngs


# ─────────────────────────── shared ───────────────────────────

def check_lesson_consistency(project: Path, slides_meta) -> None:
    """Guard against the diary-authoring drift class of bugs.

    1. An unresolved `OPEN DECISION` left in brief.md means the diary structure was
       never settled — whoever builds resolves it silently (often wrong).
    2. A diary-card with NO `teacher_note` key likely dropped its teacher layer.
       Set `"teacher_note": null` to opt a card out (e.g. a closing prayer card).
    """
    brief = project / "brief.md"
    if brief.exists() and "OPEN DECISION" in brief.read_text():
        print("[GUARD] brief.md still contains an unresolved OPEN DECISION — "
              "resolve it before trusting the diary/content structure.")
    for i, s in enumerate(slides_meta, 1):
        if s.get("type") == "diary-card" and "teacher_note" not in s:
            print(f"[GUARD] slide {i} diary-card '{s.get('title','?')}' has no "
                  f"teacher_note (set it to null to intentionally opt out).")


def build_collections(project: Path, scene_ids, resolve_image) -> None:
    """Copy every scene's final image into output/collections/<slug>.png and zip it."""
    collections_dir = project / "output" / "collections"
    if collections_dir.exists():
        shutil.rmtree(collections_dir)
    collections_dir.mkdir(parents=True, exist_ok=True)

    collected = []
    for scene_id in scene_ids:
        img_path = resolve_image(scene_id)
        if img_path is None:
            continue
        slug = scene_id.replace("/", "_")
        dest = collections_dir / f"{slug}.png"
        shutil.copy2(img_path, dest)
        collected.append(dest)
        print(f"[COLLECT] {scene_id} → collections/{dest.name}")

    if not collected:
        print("[WARN]  no images collected — nothing to zip")
        return

    zip_path = project / "output" / "collections.zip"
    with zipfile.ZipFile(zip_path, "w", zipfile.ZIP_DEFLATED) as zf:
        for f in sorted(collected):
            zf.write(f, arcname=f"collections/{f.name}")
    print(f"[ZIPPED] {len(collected)} image(s) → {zip_path}")


def make_resolver(project: Path, scenes_by_id):
    def resolve_image(slide_id: str):
        scene = scenes_by_id.get(slide_id, {})
        mode = scene.get("mode", "")
        if mode == "reuse":
            source_id = scene.get("source", "")
            source_project_name = scene.get("source_project")
            src_project = SLIDES_DIR / "projects" / source_project_name if source_project_name else project
            for fname in ("output.png", "draft.png"):
                p = src_project / "scenes" / source_id / fname
                if p.exists():
                    return p
            return None
        for fname in ("output.png", "draft.png"):
            p = project / "scenes" / slide_id / fname
            if p.exists():
                return p
        return None
    return resolve_image


def build_deck(project_name: str, out_name: str | None = None, native: bool = True) -> Path:
    project = SLIDES_DIR / "projects" / project_name
    slides_file = project / "slides.json"
    if not slides_file.exists():
        print(f"[ERROR] {slides_file} not found")
        sys.exit(1)

    slides_meta = json.loads(slides_file.read_text())

    scenes_file = project / "scenes.json"
    scenes_by_id = {}
    if scenes_file.exists():
        for s in json.loads(scenes_file.read_text()):
            scenes_by_id[s["id"]] = s
    resolve_image = make_resolver(project, scenes_by_id)

    check_lesson_consistency(project, slides_meta)

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    typed = any("type" in s for s in slides_meta)
    if typed and native:
        print(f"[MODE]  typed lesson deck — NATIVE editable text ({len(slides_meta)} slides)")
        for i, slide in enumerate(slides_meta, 1):
            deck_render.render_slide(prs, slide, resolve_image)
            print(f"[NATIVE] slide-{i:02d} ({slide.get('type')})")
    elif typed:
        print(f"[MODE]  typed lesson deck — screenshot ({len(slides_meta)} slides)")
        for png in render_typed_deck(slides_meta, resolve_image):
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            slide.shapes.add_picture(str(png), Emu(0), Emu(0), SLIDE_W, SLIDE_H)
    else:
        print(f"[MODE]  legacy image deck")
        for slide in slides_meta:
            img_path = resolve_image(slide["id"])
            if img_path is None:
                print(f"[MISSING] {slide['id']} — no image found, skipping")
                continue
            suffix = " (draft)" if img_path.name == "draft.png" else ""
            print(f"[SLIDE]  {slide['id']}{suffix} → {img_path.name}")
            add_image_slide(prs, img_path, slide.get("theme", ""))
            if slide.get("cta"):
                add_cta_slide(prs, slide["cta"], slide.get("theme", ""))

    out_path = project / "output" / (out_name or f"{project_name}.pptx")
    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out_path))
    print(f"[SAVED] {out_path}")

    build_collections(project, list(scenes_by_id.keys()), resolve_image)
    return out_path


def main():
    args = sys.argv[1:]
    if not args:
        print("Usage: python build_pptx.py <project> [--out filename.pptx]")
        sys.exit(1)
    project_name = args[0]
    out_name = args[args.index("--out") + 1] if "--out" in args else None
    # native (editable text) is the default for typed decks; --screenshot opts into
    # the pixel-perfect flat-image path instead.
    native = "--screenshot" not in args
    build_deck(project_name, out_name, native=native)


if __name__ == "__main__":
    main()
