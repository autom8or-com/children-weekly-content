"""Native python-pptx renderer for the storybook lesson deck.

Produces EDITABLE PowerPoint text (real text boxes) + the illustration placed as
a picture — so the deck can be edited in PowerPoint / Canva afterward. Reuses the
palette + schema from deck_html.py; the look is close to the Claude Design mock
(exact palette/layout/fonts; cards have uniform corners and no divider bar).

Fonts Caveat + Patrick Hand must be installed on whatever machine opens the file.
"""
import hashlib
import re
import tempfile
from pathlib import Path

from PIL import Image
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

from deck_html import PALETTES, BG_CREAM, BG_DARK, TEXT_DARK, TEXT_BROWN, resolve_date

FONT_HEAD = "Caveat"
FONT_LABEL = "Patrick Hand"
IN_PER_PX = 13.33 / 1920

# the illustration only displays at ~1200px on the 1920px stage, so embedding the
# full 2K source bloats the .pptx (~5MB each). Downscale + recompress on the way in;
# originals are left untouched (collections still uses full-res).
IMG_MAX_W = 1600
IMG_QUALITY = 88
_IMG_CACHE = Path(tempfile.mkdtemp(prefix="deck-img-"))


def compress_image(path):
    """Return a downscaled, recompressed copy of `path` (cached by content+mtime)."""
    src = Path(path)
    try:
        mtime = src.stat().st_mtime
    except OSError:
        return path
    key = hashlib.md5(f"{src}|{mtime}|{IMG_MAX_W}|{IMG_QUALITY}".encode()).hexdigest()[:16]
    out = _IMG_CACHE / f"{key}.jpg"
    if out.exists():
        return out
    try:
        im = Image.open(src).convert("RGB")
    except Exception:
        return path  # unreadable — fall back to the original
    if im.width > IMG_MAX_W:
        im = im.resize((IMG_MAX_W, round(im.height * IMG_MAX_W / im.width)), Image.LANCZOS)
    im.save(out, "JPEG", quality=IMG_QUALITY, optimize=True)
    return out

# Card geometry (px on the 1920x1080 stage) — extends below the image so longer
# Pre-K story / diary bodies have room. Image keeps 16:9 at 675 px. CARD_H must
# leave room for the footer (footer is pinned ~30 px from the card bottom).
CARD_X, TOP, CARD_W, CARD_H = 90, 202.5, 540, 850
# Card interior padding (matches HTML: 40 top/bottom, 48 sides, 18 row gap).
PAD_X, PAD_Y = 44, 28
# Illustration keeps its 16:9 aspect at the original height; card extends taller.
IMG_X, IMG_W, _ILLUST_H = 690, 1140, 675
# Teacher's-note band (matches TN_BAND/TN_GAP if added later for diary-cards).
TN_BAND, TN_GAP = 150, 16

# Layout constants for the body textbox inside a card. These keep the textbox
# strictly within the card so it can never bleed past the card edge.
CARD_INNER_LEFT_PX = CARD_X + PAD_X
CARD_INNER_TOP_PX = TOP + PAD_Y
CARD_INNER_W_PX = CARD_W - 2 * PAD_X
# Body textbox: leave ~38 px at the bottom for the footer (18pt + padding).
CARD_INNER_BODY_H_PX = CARD_H - 2 * PAD_Y - 38
CARD_FOOTER_TOP_PX = TOP + CARD_H - PAD_Y - 18
CARD_FOOTER_H_PX = 18


def px(n):
    return Inches(n * IN_PER_PX)


def C(hex6):
    return RGBColor.from_string(hex6)


def _bg(slide, hex6):
    f = slide.background.fill
    f.solid()
    f.fore_color.rgb = C(hex6)


def _ring(slide, inset_px, width_px, hex6):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, px(inset_px), px(inset_px),
                                 px(1920 - 2 * inset_px), px(1080 - 2 * inset_px))
    shp.fill.background()
    shp.line.color.rgb = C(hex6)
    shp.line.width = px(width_px)
    shp.shadow.inherit = False


def _rect(slide, x, y, w, h, hex6, rounded=False):
    shape = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
    shp = slide.shapes.add_shape(shape, px(x), px(y), px(w), px(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = C(hex6)
    shp.line.fill.background()
    shp.shadow.inherit = False
    return shp


def _runs(paragraph, text, *, font, size, color, bold):
    """Add runs to a paragraph, turning **bold** segments bold."""
    if bold is None:
        bold = False
    for i, seg in enumerate(str(text).split("**")):
        if seg == "":
            continue
        run = paragraph.add_run()
        run.text = seg
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold or (i % 2 == 1)
        run.font.color.rgb = C(color)


# Calibrated for Caveat / Patrick Hand in PowerPoint/LibreOffice (Caveat renders
# ~1.5x wider here than in Chrome, so we measure conservatively to keep text in-card).
# The auto-shrink logic depends on these to know when to scale down. We deliberately
# overestimate char width so the auto-shrink errs on the side of more shrinkage.
_CAVEAT_CHAR_IN = {16: 0.115, 18: 0.13, 19: 0.135, 20: 0.14, 22: 0.155, 24: 0.17,
                   26: 0.18, 28: 0.19, 30: 0.205, 32: 0.22, 40: 0.27, 42: 0.285,
                   50: 0.34, 52: 0.35, 56: 0.375, 60: 0.40, 76: 0.51, 80: 0.535,
                   84: 0.56, 120: 0.80, 132: 0.88, 140: 0.93, 148: 0.98}
_PATRICK_CHAR_IN = {20: 0.10, 22: 0.11, 24: 0.12, 26: 0.13, 28: 0.14, 30: 0.15,
                    32: 0.16, 34: 0.17, 40: 0.20, 42: 0.21}


def _char_width_in(font, size):
    """Approximate avg character width in inches for the given font + size."""
    if font == FONT_HEAD:
        table = _CAVEAT_CHAR_IN
    else:
        table = _PATRICK_CHAR_IN
    # pick closest table entry
    if size in table:
        return table[size]
    keys = sorted(table.keys())
    if size <= keys[0]:
        return table[keys[0]] * size / keys[0]
    if size >= keys[-1]:
        return table[keys[-1]] * size / keys[-1]
    # linear interp
    for k0, k1 in zip(keys, keys[1:]):
        if k0 <= size <= k1:
            r = (size - k0) / (k1 - k0)
            return table[k0] * (1 - r) + table[k1] * r
    return table[keys[-1]]


def _lines_for_text(text, width_in, font, size):
    """Estimate number of wrapped lines for `text` at given width / font / size."""
    cw = _char_width_in(font, size)
    chars_per_line = max(5, int(width_in / cw))
    n = 0
    for paragraph in str(text).split("\n"):
        if not paragraph.strip():
            n += 1
            continue
        # word-wrap
        line_count = 0
        cur = 0
        for word in paragraph.split():
            wlen = len(word) + 1
            if cur == 0:
                cur = wlen
            elif cur + wlen > chars_per_line:
                line_count += 1
                cur = wlen
            else:
                cur += wlen
        n += line_count + 1
    return n


def _pt_height(size, line_spacing=1.6):
    """Height in points for one line at given font size + line spacing."""
    return size * line_spacing


def _textbox(slide, x, y, w, h, lines, *, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
             auto_fit_shrink=True):
    """Add a textbox. If `auto_fit_shrink` and the text would overflow `h`, shrink
    body sizes by 1 pt at a time until it fits (or we hit a floor of 18 pt).
    All dimensions (x, y, w, h) are in INCHES."""
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0

    # Auto-shrink: scale down `size` for body paragraphs (skip title/chip if they
    # fit; only shrink everything proportionally as a last resort).
    if auto_fit_shrink and h > 0:
        # Total available height in points (1 in = 72 pt)
        avail_pt = h * 72
        # Estimate total pt used (sum of line count * line height)
        def used_pt(scale=1.0):
            total = 0.0
            for ln in lines:
                for sub in str(ln["text"]).split("\n"):
                    n_lines = _lines_for_text(sub, w, ln.get("font", FONT_HEAD),
                                              max(10, int(ln["size"] * scale)))
                    total += n_lines * _pt_height(ln["size"] * scale, ln.get("line_spacing", 1.1))
                    # paragraph space_after in pt
                    sa = ln.get("space_after", 0)
                    if sa:
                        total += sa
            return total
        scale = 1.0
        # Try 1.0 first; if it overflows, shrink in 0.05 steps.
        if used_pt(1.0) > avail_pt:
            for s in [0.95, 0.90, 0.85, 0.80, 0.75, 0.70, 0.65]:
                if used_pt(s) <= avail_pt:
                    scale = s
                    break
            else:
                scale = 0.65  # floor — overflow is the lesser evil vs illegible text

    first = True
    for ln in lines:
        size = max(10, int(ln["size"] * scale)) if auto_fit_shrink else ln["size"]
        for sub in str(ln["text"]).split("\n"):
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False
            p.alignment = align
            p.space_after = Pt(ln.get("space_after", 6))
            p.line_spacing = ln.get("line_spacing", 1.1)
            _runs(p, sub, font=ln.get("font", FONT_HEAD), size=size,
                  color=ln["color"], bold=ln.get("bold", None))
    return tb


def render_bookend(prs, *, bg, frame, lines):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide, bg)
    if frame:
        _ring(slide, 18, 16, frame)
        _ring(slide, 46, 5, frame)
    _textbox(slide, 160 * IN_PER_PX, 0, 1600 * IN_PER_PX, 1080 * IN_PER_PX, lines,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, auto_fit_shrink=False)
    return slide


def render_card(prs, *, sec_bg, card_bg, accent, chip, chip_font, chip_col,
                title, body_lines, footer, footer_col, image_path, stripe=None,
                teacher_note=None, footer_size=16, body_size=16, chip_size=None,
                title_size=36):
    """Render the standard left-card + right-image slide. body_lines is a list of
    dicts (chip + title rows are auto-prepended)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide, sec_bg)

    if teacher_note:
        card_h = CARD_H - (TN_BAND + TN_GAP)
    else:
        card_h = CARD_H
    illust_h = _ILLUST_H

    # Body textbox dimensions — must be computed from the ACTUAL card height so
    # the text never bleeds past the card edge (esp. for diary-cards that share
    # the stage with a teacher's-note band). Convert to inches for _textbox.
    body_h_in = (card_h - 2 * PAD_Y - 38) * IN_PER_PX
    footer_top_in = (TOP + card_h - PAD_Y - 18) * IN_PER_PX
    inner_left_in = (CARD_X + PAD_X) * IN_PER_PX
    inner_top_in = (TOP + PAD_Y) * IN_PER_PX
    inner_w_in = (CARD_W - 2 * PAD_X) * IN_PER_PX
    footer_h_in = CARD_FOOTER_H_PX * IN_PER_PX

    # illustration (right)
    if image_path:
        slide.shapes.add_picture(str(compress_image(image_path)), px(IMG_X), px(TOP), px(IMG_W), px(illust_h))
    else:
        _rect(slide, IMG_X, TOP, IMG_W, illust_h, "D8C9A8", rounded=True)

    # card (left)
    _rect(slide, CARD_X, TOP, CARD_W, card_h, card_bg, rounded=True)
    if stripe:
        _rect(slide, CARD_X, TOP, 8, card_h, stripe)

    # body textbox — strict card interior, header (chip + title) on top, then body.
    if chip_size is None:
        chip_size = 16 if chip_font == FONT_LABEL else 18
    head = [
        {"text": chip, "font": chip_font, "size": chip_size, "color": chip_col,
         "space_after": 6, "line_spacing": 1.0},
        {"text": title, "size": title_size, "color": TEXT_DARK, "bold": True,
         "space_after": 4, "line_spacing": 1.0},
        # accent divider
        {"text": "▬", "font": FONT_LABEL, "size": 10, "color": accent,
         "space_after": 4, "line_spacing": 1.0},
    ]
    head += body_lines
    # Tighten the body textbox to leave room for the footer at the bottom of the card.
    # _textbox takes dimensions in INCHES (px() converts to EMU internally).
    _textbox(slide, inner_left_in, inner_top_in, inner_w_in, body_h_in, head)
    # footer (prompt / signature) — pinned to the bottom interior of the card.
    if footer:
        _textbox(slide, inner_left_in, footer_top_in, inner_w_in, footer_h_in,
                 [{"text": footer, "size": footer_size, "color": footer_col,
                   "space_after": 0, "line_spacing": 1.0}],
                 auto_fit_shrink=False)

    # teacher's note band
    if teacher_note:
        ty = TOP + card_h + TN_GAP
        tw = (IMG_X + IMG_W) - CARD_X
        _rect(slide, CARD_X, ty, tw, TN_BAND, "ECE3F8", rounded=True)
        _rect(slide, CARD_X, ty, 8, TN_BAND, "7B4DB5")
        # Convert to inches for _textbox.
        _textbox(slide, (CARD_X + 28) * IN_PER_PX, (ty + 16) * IN_PER_PX,
                 (tw - 56) * IN_PER_PX, (TN_BAND - 28) * IN_PER_PX,
                 [{"text": "📝  Teacher's Note", "font": FONT_LABEL, "size": 20,
                   "color": "7B4DB5", "space_after": 6, "line_spacing": 1.0},
                  {"text": teacher_note, "size": 19, "color": TEXT_DARK,
                   "space_after": 0, "line_spacing": 1.1}],
                 auto_fit_shrink=True)
    return slide


def _body_lines(body, body_size=16):
    """Caveat renders wider in PowerPoint/LibreOffice than in Chrome; 16pt is the
    default body size — long Pre-K story/diary cards shrink further via auto-fit."""
    return [{"text": p["text"], "size": body_size,
             "color": (TEXT_BROWN if p.get("strong") else TEXT_DARK),
             "space_after": 4, "line_spacing": 1.3} for p in body]


def _bullet_lines(bullets, body_size=16):
    return [{"text": f'{b.get("emoji","•")}  {b["text"]}', "size": body_size, "color": TEXT_DARK,
             "space_after": 6, "line_spacing": 1.25} for b in bullets]


def _chip_text(s):
    """Extract emoji + first non-emoji chunk for chip rendering, preserving emoji."""
    # Used by the section-header etc. — pass through as-is.
    return s


# ─────────────────────────── bookend variants ───────────────────────────

def render_topic(prs, s, image):
    """topic — single dark page with eyebrow + huge title (optional image below)."""
    pal = PALETTES["preamble"]  # color isn't used here; fixed scheme
    lines = [
        {"text": s.get("eyebrow", "✦ TOPIC ✦"), "font": FONT_LABEL, "size": 30,
         "color": "C97B3A", "space_after": 24, "line_spacing": 1.0},
        {"text": s["title"], "size": 96, "color": "FDF6E3", "bold": True,
         "space_after": 0, "line_spacing": 0.95},
    ]
    return render_bookend(prs, bg=BG_DARK, frame="C97B3A", lines=lines)


def render_verse(prs, s, image):
    """verse — dark page with quote + reference (matches deck_html _bookend verse)."""
    lines = [
        {"text": s["eyebrow"], "font": FONT_LABEL, "size": 26, "color": "C97B3A",
         "space_after": 24, "line_spacing": 1.0},
        {"text": s["title"], "size": 56, "color": "FDF6E3", "bold": True,
         "space_after": 20, "line_spacing": 1.1},
        {"text": "▬▬▬", "font": FONT_LABEL, "size": 16, "color": "C97B3A",
         "space_after": 18, "line_spacing": 1.0},
        {"text": s["reference"], "font": FONT_LABEL, "size": 30, "color": "EDD9A3",
         "space_after": 0, "line_spacing": 1.0},
    ]
    return render_bookend(prs, bg=BG_DARK, frame="C97B3A", lines=lines)


def render_prayer(prs, s, image):
    lines = [
        {"text": s["eyebrow"], "font": FONT_LABEL, "size": 26, "color": "C97B3A",
         "space_after": 24, "line_spacing": 1.0},
        {"text": s["title"], "size": 44, "color": TEXT_DARK, "bold": False,
         "space_after": 18, "line_spacing": 1.2},
    ]
    if s.get("footer"):
        lines.append({"text": s["footer"], "font": FONT_LABEL, "size": 24,
                      "color": TEXT_BROWN, "space_after": 0, "line_spacing": 1.0})
    return render_bookend(prs, bg=BG_CREAM, frame="EDD9A3", lines=lines)


def render_title_goodbye(prs, s, is_goodbye):
    dark = is_goodbye
    bg = BG_DARK if dark else BG_CREAM
    title_col = "FDF6E3" if dark else TEXT_DARK
    sub_col = "EDD9A3" if dark else TEXT_BROWN
    lines = [
        {"text": s["eyebrow"], "font": FONT_LABEL, "size": 28, "color": "C97B3A",
         "space_after": 20, "line_spacing": 1.0},
        {"text": s["title"], "size": 96, "color": title_col, "bold": True,
         "space_after": 20, "line_spacing": 0.95},
        {"text": "▬▬▬", "font": FONT_LABEL, "size": 18, "color": "C97B3A",
         "space_after": 16, "line_spacing": 1.0},
    ]
    if s.get("subtitle"):
        lines.append({"text": s["subtitle"], "size": 38, "color": sub_col,
                      "space_after": 12, "line_spacing": 1.0})
    if s.get("date"):
        lines.append({"text": resolve_date(s["date"]), "font": FONT_LABEL, "size": 26,
                      "color": "C97B3A", "space_after": 0, "line_spacing": 1.0})
    if s.get("footer"):
        lines.append({"text": s["footer"], "font": FONT_LABEL, "size": 22,
                      "color": "A0856A" if not dark else "EDD9A3",
                      "space_after": 0, "line_spacing": 1.0})
    return render_bookend(prs, bg=bg, frame="C97B3A", lines=lines)


def render_section_header(prs, s):
    pal = PALETTES.get(s["palette"], PALETTES["preamble"])
    lines = [
        {"text": s["eyebrow"], "font": FONT_LABEL, "size": 30, "color": pal["hdr_text"],
         "space_after": 24, "line_spacing": 1.0},
        {"text": s["title"], "size": 110, "color": pal["hdr_title"], "bold": True,
         "space_after": 20, "line_spacing": 0.95},
        {"text": "▬▬▬", "font": FONT_LABEL, "size": 18, "color": pal["hdr_title"],
         "space_after": 20, "line_spacing": 1.0},
        {"text": s["footer"], "size": 38, "color": pal["hdr_text"],
         "space_after": 0, "line_spacing": 1.0},
    ]
    return render_bookend(prs, bg=pal["hdr_bg"], frame=None, lines=lines)


# ─────────────────────────── card variants ───────────────────────────

def render_objectives(prs, s, image):
    """objectives — same shape as story-card but with lead + bullets."""
    pal = PALETTES.get(s.get("palette", "preamble"), PALETTES["preamble"])
    head_extras = []
    if s.get("lead"):
        head_extras.append({"text": s["lead"], "size": 18, "color": TEXT_BROWN,
                            "space_after": 4, "line_spacing": 1.2})
    body = head_extras + _bullet_lines(s.get("bullets", []), body_size=18)
    return render_card(prs, sec_bg=pal["sec_bg"], card_bg=pal["card"], accent=pal["accent"],
                       chip=s.get("chip", "🎯 Objectives"), chip_font=FONT_LABEL,
                       chip_col=pal["accent"], title=s.get("title", "Objectives"),
                       body_lines=body, footer=s.get("prompt", ""), footer_col=TEXT_BROWN,
                       image_path=image, body_size=18, chip_size=18, title_size=40)


def render_house_rules(prs, s):
    pal = PALETTES.get(s.get("palette", "preamble"), PALETTES["preamble"])
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide, pal["sec_bg"])

    # Header (eyebrow + title) — top centre
    head = [
        {"text": s.get("chip", "📋 Before We Begin"), "font": FONT_LABEL, "size": 32,
         "color": pal["accent"], "space_after": 16, "line_spacing": 1.0},
        {"text": s.get("title", "House Rules"), "size": 80, "color": pal["hdr_title"],
         "bold": True, "space_after": 0, "line_spacing": 1.0},
    ]
    _textbox(slide, 240 * IN_PER_PX, 100 * IN_PER_PX, 1440 * IN_PER_PX, 220 * IN_PER_PX,
             head, align=PP_ALIGN.CENTER, auto_fit_shrink=False)

    # Rules rows — stacked, each in a rounded rectangle
    rules = s.get("rules", [])
    n = len(rules)
    row_h = 100
    gap = 22
    total_h = n * row_h + (n - 1) * gap
    y0 = 380
    for i, r in enumerate(rules):
        ry = y0 + i * (row_h + gap)
        _rect(slide, 360, ry, 1200, row_h, pal["card"], rounded=True)
        # icon
        icon_text = r.get("emoji", "•")
        _textbox(slide, 388 * IN_PER_PX, (ry + 18) * IN_PER_PX, 80 * IN_PER_PX, (row_h - 36) * IN_PER_PX,
                 [{"text": icon_text, "size": 44, "color": TEXT_DARK, "space_after": 0}],
                 align=PP_ALIGN.CENTER, auto_fit_shrink=False)
        # text
        _textbox(slide, 480 * IN_PER_PX, (ry + 16) * IN_PER_PX, 1060 * IN_PER_PX, (row_h - 32) * IN_PER_PX,
                 [{"text": r["text"], "size": 28, "color": TEXT_DARK, "space_after": 0,
                   "line_spacing": 1.2}], auto_fit_shrink=True)
    return slide


def render_outlines(prs, s, resolve_image):
    pal = PALETTES.get(s.get("palette", "preamble"), PALETTES["preamble"])
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide, pal["sec_bg"])

    # Header
    head = [
        {"text": s.get("chip", "🗺️ Today's Journey"), "font": FONT_LABEL, "size": 32,
         "color": pal["accent"], "space_after": 16, "line_spacing": 1.0},
        {"text": s.get("title", "Outlines"), "size": 84, "color": pal["hdr_title"],
         "bold": True, "space_after": 0, "line_spacing": 1.0},
    ]
    _textbox(slide, 240 * IN_PER_PX, 100 * IN_PER_PX, 1440 * IN_PER_PX, 220 * IN_PER_PX,
             head, align=PP_ALIGN.CENTER, auto_fit_shrink=False)

    # Items: 4 columns of (image + label)
    items = s.get("items", [])
    n = len(items)
    if n == 0:
        return slide
    # layout: equal columns centred, with gap
    col_w = 360
    gap = 40
    total_w = n * col_w + (n - 1) * gap
    x0 = (1920 - total_w) / 2
    y = 420
    for i, it in enumerate(items):
        x = x0 + i * (col_w + gap)
        # thumb
        img = resolve_image(it.get("image")) if it.get("image") else None
        if img:
            slide.shapes.add_picture(str(compress_image(img)), px(x), px(y), px(col_w), px(300))
        else:
            _rect(slide, x, y, col_w, 300, pal["card"], rounded=True)
        # label
        _textbox(slide, x * IN_PER_PX, (y + 320) * IN_PER_PX, col_w * IN_PER_PX, 80 * IN_PER_PX,
                 [{"text": f"{i+1}. {it['label']}", "size": 32, "color": pal["hdr_title"],
                   "bold": True, "space_after": 0, "line_spacing": 1.0}],
                 align=PP_ALIGN.CENTER, auto_fit_shrink=True)
    return slide


def render_bible_text(prs, s):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide, BG_CREAM)
    _ring(slide, 18, 16, "C97B3A")  # outer orange frame

    body = s.get("body", [])
    cols = s.get("columns", 2 if len(body) > 1 else 1)

    # Header
    head = [
        {"text": s.get("eyebrow", "📖 Bible Text"), "font": FONT_LABEL, "size": 26,
         "color": "C97B3A", "space_after": 8, "line_spacing": 1.0},
        {"text": s["reference"], "size": 44, "color": TEXT_DARK, "bold": True,
         "space_after": 16, "line_spacing": 1.0},
    ]
    _textbox(slide, 110 * IN_PER_PX, 60 * IN_PER_PX, 1700 * IN_PER_PX, 240 * IN_PER_PX,
             head, auto_fit_shrink=True)

    # Body — single or two-column layout
    if cols == 1:
        body_lines = [{"text": p, "size": 22, "color": TEXT_DARK, "space_after": 14,
                       "line_spacing": 1.5} for p in body]
        _textbox(slide, 110 * IN_PER_PX, 320 * IN_PER_PX, 1700 * IN_PER_PX, 720 * IN_PER_PX,
                 body_lines, auto_fit_shrink=True)
    else:
        # Distribute paragraphs between two columns; render each column
        # separately so column-count CSS doesn't fight PowerPoint layout.
        mid = (len(body) + 1) // 2
        col1 = body[:mid]
        col2 = body[mid:]
        col_lines_1 = [{"text": p, "size": 22, "color": TEXT_DARK, "space_after": 14,
                        "line_spacing": 1.5} for p in col1]
        col_lines_2 = [{"text": p, "size": 22, "color": TEXT_DARK, "space_after": 14,
                        "line_spacing": 1.5} for p in col2]
        _textbox(slide, 110 * IN_PER_PX, 320 * IN_PER_PX, 800 * IN_PER_PX, 700 * IN_PER_PX,
                 col_lines_1, auto_fit_shrink=True)
        _textbox(slide, 970 * IN_PER_PX, 320 * IN_PER_PX, 800 * IN_PER_PX, 700 * IN_PER_PX,
                 col_lines_2, auto_fit_shrink=True)
    return slide


def render_application(prs, s, image):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide, BG_CREAM)
    _ring(slide, 18, 16, "C97B3A")

    # Left column: text, right column: image
    left = [
        {"text": s.get("title", "Saved by the blood of the Lamb."), "size": 44,
         "color": TEXT_DARK, "bold": True, "space_after": 16, "line_spacing": 1.0},
    ]
    for line in s.get("intro", []):
        left.append({"text": line, "size": 26, "color": TEXT_BROWN,
                     "space_after": 10, "line_spacing": 1.3})
    left.append({"text": s.get("steps_lead", "To be rescued:"), "font": FONT_LABEL,
                 "size": 28, "color": "C97B3A", "space_after": 12, "line_spacing": 1.0})
    # Default A-B-C steps if not provided
    from deck_html import ABC_STEPS
    for st in s.get("steps", ABC_STEPS):
        left.append({"text": f"{st['letter']} – {st['text']}", "size": 30,
                     "color": TEXT_DARK, "bold": True, "space_after": 8, "line_spacing": 1.2})
    _textbox(slide, 120 * IN_PER_PX, 110 * IN_PER_PX, 1180 * IN_PER_PX, 880 * IN_PER_PX,
             left, auto_fit_shrink=True)

    # Image on the right
    if image:
        slide.shapes.add_picture(str(compress_image(image)), px(1380), px(220), px(420), px(660))
    return slide


# ─────────────────────────── dispatch ───────────────────────────

def render_slide(prs, s, resolve_image):
    """Render one typed slide dict natively. `s['image']` is a scene id or None."""
    t = s.get("type")
    image = None
    if s.get("image"):
        image = resolve_image(s["image"])
    if t == "section-header" and s.get("thumbs"):
        return render_section_header_with_thumbs(prs, s, resolve_image)

    if t == "title":
        return render_title_goodbye(prs, s, is_goodbye=False)
    if t == "goodbye":
        return render_title_goodbye(prs, s, is_goodbye=True)
    if t == "topic":
        return render_topic(prs, s, image)
    if t == "section-header":
        return render_section_header(prs, s)
    if t == "verse":
        return render_verse(prs, s, image)
    if t == "prayer":
        return render_prayer(prs, s, image)

    if t == "house-rules":
        return render_house_rules(prs, s)
    if t == "outlines":
        return render_outlines(prs, s, resolve_image)
    if t == "bible-text":
        return render_bible_text(prs, s)
    if t == "application":
        return render_application(prs, s, image)
    if t == "objectives":
        return render_objectives(prs, s, image)

    if t == "story-card":
        pal = PALETTES.get(s.get("palette", "preamble"), PALETTES["flies"])
        return render_card(prs, sec_bg=pal["sec_bg"], card_bg=pal["card"], accent=pal["accent"],
                           chip=s["chip"], chip_font=FONT_LABEL, chip_col=pal["accent"],
                           title=s["title"], body_lines=_body_lines(s.get("body", [])),
                           footer=s.get("prompt", ""), footer_col=TEXT_BROWN, image_path=image)

    if t == "diary-card":
        pal = PALETTES["diary"]
        return render_card(prs, sec_bg=pal["sec_bg"], card_bg=pal["card"], accent=pal["accent"],
                           chip=s["chip"], chip_font=FONT_HEAD, chip_col="7B4DB5",
                           title=s["title"], body_lines=_body_lines(s.get("body", [])),
                           footer=s.get("signature", ""), footer_col="7B4DB5",
                           image_path=image, stripe=pal["accent"],
                           teacher_note=s.get("teacher_note"))

    if t == "summary":
        pal = PALETTES.get(s.get("palette", "preamble"), PALETTES["preamble"])
        return render_card(prs, sec_bg=pal["sec_bg"], card_bg=pal["card"], accent=pal["accent"],
                           chip=s["chip"], chip_font=FONT_LABEL, chip_col=pal["accent"],
                           title=s["title"], body_lines=_bullet_lines(s.get("bullets", [])),
                           footer=s.get("prompt", ""), footer_col=TEXT_BROWN, image_path=image)

    raise ValueError(f"unknown slide type: {t!r}")


def render_section_header_with_thumbs(prs, s, resolve_image):
    pal = PALETTES.get(s["palette"], PALETTES["preamble"])
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide, pal["hdr_bg"])

    lines = [
        {"text": s["eyebrow"], "font": FONT_LABEL, "size": 30, "color": pal["hdr_text"],
         "space_after": 24, "line_spacing": 1.0},
        {"text": s["title"], "size": 100, "color": pal["hdr_title"], "bold": True,
         "space_after": 16, "line_spacing": 0.95},
        {"text": "▬▬▬", "font": FONT_LABEL, "size": 18, "color": pal["hdr_title"],
         "space_after": 18, "line_spacing": 1.0},
        {"text": s["footer"], "size": 32, "color": pal["hdr_text"],
         "space_after": 0, "line_spacing": 1.0},
    ]
    _textbox(slide, 160 * IN_PER_PX, 150 * IN_PER_PX, 1600 * IN_PER_PX, 480 * IN_PER_PX,
             lines, align=PP_ALIGN.CENTER, auto_fit_shrink=True)

    # Thumbs at the bottom
    thumbs = [resolve_image(t) for t in s.get("thumbs", []) if t]
    thumbs = [t for t in thumbs if t]
    n = len(thumbs)
    if n:
        tw, th = 300, 200
        gap = 40
        total_w = n * tw + (n - 1) * gap
        x0 = (1920 - total_w) / 2
        y = 800
        for i, t in enumerate(thumbs):
            x = x0 + i * (tw + gap)
            slide.shapes.add_picture(str(compress_image(t)), px(x), px(y), px(tw), px(th))
    return slide
